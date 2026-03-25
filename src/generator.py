"""
生成模块
使用 python-pptx 基于结构化数据生成 PPTX 文件。
"""
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

class GeneratorError(Exception):
    pass

class PPTGenerator:
    """使用 python-pptx 生成 PPT"""

    def __init__(self, config: dict):
        self.output_dir = config.get("ppt", {}).get("output_dir", "./output")
        self.font_config = config.get("ppt", {}).get("font", {})
        
        # 确保输出目录存在
        os.makedirs(self.output_dir, exist_ok=True)

    def generate(self, slides_data: list, title_text: str = "NotebookLM AI 生成 PPT") -> str:
        """生成并保存 PPTX 文件"""
        print("[Generator] 开始生成 PPT...")
        prs = Presentation()
        
        # 设置页面大小 16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # PPTX 默认版式 (0=标题幻灯片, 1=标题和内容)
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]

        # 封面页
        print("[Generator] -> 添加封面页")
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        subtitle.text = "Powered by NotebookLM & Playwright"

        # 应用字体
        self._apply_font(title.text_frame, title=True)
        self._apply_font(subtitle.text_frame, title=False)

        # 内容页
        for i, slide_data in enumerate(slides_data):
            print(f"[Generator] -> 添加内容页 {i+1}: {slide_data.get('title')}")
            slide = prs.slides.add_slide(content_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]

            # 填充标题
            title_shape.text = slide_data.get("title", f"Slide {i+1}")
            self._apply_font(title_shape.text_frame, title=True)

            # 填充要点
            text_frame = body_shape.text_frame
            text_frame.clear()  # 清除默认文本
            
            bullets = slide_data.get("bullets", [])
            for bullet_text in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet_text
                p.space_after = Pt(14)
            
            self._apply_font(text_frame, title=False)

            # 填充演讲备注
            notes = slide_data.get("notes", "")
            if notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes

        # 保存文件
        output_filename = f"{title_text.replace(' ', '_')}.pptx"
        output_path = os.path.join(self.output_dir, output_filename)
        
        try:
            prs.save(output_path)
        except Exception as e:
            raise GeneratorError(f"保存 PPT 失败: {str(e)}")
        
        print(f"[Generator] [OK] PPT 生成完毕! 保存在: {os.path.abspath(output_path)}")
        return output_path

    def _apply_font(self, text_frame, title=False):
        """简单封装：为段落应用统一定义的字体。"""
        font_name = self.font_config.get("title" if title else "body", "Arial")
        for paragraph in text_frame.paragraphs:
             for run in paragraph.runs:
                 run.font.name = font_name
                 # MVP 简单设置字体大小
                 if title:
                     run.font.size = Pt(44)
                 else:
                     run.font.size = Pt(28)
